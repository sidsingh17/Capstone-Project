"""
AI-Powered Supply Chain Risk Intelligence Assistant
Presentation Generator — Prodapt Solutions Brand Theme
Navy #0F1F3E / #1C3461   Orange #F05A28
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

# ── Brand Colors ───────────────────────────────────────────────────────────
NAVY_BG    = RGBColor(0x0F, 0x1F, 0x3E)
NAVY       = RGBColor(0x1C, 0x34, 0x61)
NAVY_LIGHT = RGBColor(0x24, 0x3D, 0x72)
ORANGE     = RGBColor(0xF0, 0x5A, 0x28)
ORANGE_LT  = RGBColor(0xFE, 0xF0, 0xEB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG    = RGBColor(0xF4, 0xF6, 0xFB)
CARD_BG    = RGBColor(0xEE, 0xF2, 0xF8)
TEXT       = RGBColor(0x1A, 0x20, 0x38)
MUTED      = RGBColor(0x6B, 0x7A, 0x9A)
SUCCESS    = RGBColor(0x0A, 0x7A, 0x4A)
SUCCESS_BG = RGBColor(0xE8, 0xF5, 0xEE)
WARNING    = RGBColor(0xC4, 0x7A, 0x00)
DANGER     = RGBColor(0xC8, 0x20, 0x1A)

W = Inches(13.33)   # Slide width  (16:9)
H = Inches(7.5)     # Slide height (16:9)

# ── Helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE.RECTANGLE=1
    fill_fmt = shape.fill
    if fill:
        fill_fmt.solid()
        fill_fmt.fore_color.rgb = fill
    else:
        fill_fmt.background()
    ln = shape.line
    if line:
        ln.color.rgb = line
        ln.width = line_w
    else:
        ln.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h, size=Pt(12), bold=False,
                 color=TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_label_box(slide, text, x, y, w, h, bg=NAVY, fg=WHITE, size=Pt(11),
                  bold=True, align=PP_ALIGN.CENTER, radius=False):
    box = add_rect(slide, x, y, w, h, fill=bg)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    from pptx.util import Pt as _Pt
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = fg
    return box

def arrow_right(slide, x, y, length=Inches(0.5), color=NAVY):
    line = slide.shapes.add_connector(1, x, y, x + length, y)
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    return line

def bullet_slide_body(slide, items, x, y, w, h, bullet="●",
                      size=Pt(13), color=TEXT, spacing=Inches(0.42)):
    for i, item in enumerate(items):
        add_text_box(slide, f"{bullet}  {item}",
                     x, y + i * spacing, w, Inches(0.4),
                     size=size, color=color)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def orange_line(slide, y=Inches(7.2)):
    add_rect(slide, Inches(0), y, W, Inches(0.06), fill=ORANGE)

def page_header(slide, title, subtitle=None, dark=False):
    bg_c = NAVY_BG if dark else PAGE_BG
    set_bg(slide, bg_c)
    # Top accent bar
    add_rect(slide, 0, 0, W, Inches(0.07), fill=ORANGE)
    # Header band
    add_rect(slide, 0, Inches(0.07), W, Inches(1.08), fill=NAVY)
    add_text_box(slide, title,
                 Inches(0.5), Inches(0.12), Inches(12), Inches(0.7),
                 size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.5), Inches(0.8), Inches(12), Inches(0.36),
                     size=Pt(12), color=RGBColor(0xC8, 0xD3, 0xE8),
                     align=PP_ALIGN.LEFT)
    orange_line(slide)

def card(slide, x, y, w, h, bg=CARD_BG, border=NAVY):
    add_rect(slide, x, y, w, h, fill=bg, line=border, line_w=Pt(1.2))

def prs_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = prs_slide(prs)
    set_bg(sl, NAVY_BG)
    # Orange bars
    add_rect(sl, 0, 0, W, Inches(0.12), fill=ORANGE)
    add_rect(sl, 0, Inches(7.38), W, Inches(0.12), fill=ORANGE)
    # Left accent stripe
    add_rect(sl, 0, Inches(0.12), Inches(0.1), Inches(7.26), fill=ORANGE)

    # Hexagon logo
    add_text_box(sl, "⬡", Inches(1.0), Inches(1.2), Inches(1.2), Inches(1.2),
                 size=Pt(52), color=ORANGE, align=PP_ALIGN.CENTER)

    # Company name
    add_text_box(sl, "PRODAPT SOLUTIONS",
                 Inches(2.0), Inches(1.25), Inches(9), Inches(0.55),
                 size=Pt(14), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    add_text_box(sl, "Capstone Project  |  2026",
                 Inches(2.0), Inches(1.78), Inches(9), Inches(0.4),
                 size=Pt(11), color=RGBColor(0x9A, 0xA3, 0xBB), align=PP_ALIGN.LEFT)

    # Divider
    add_rect(sl, Inches(1.8), Inches(2.3), Inches(9.5), Inches(0.03), fill=ORANGE)

    # Main title
    add_text_box(sl, "AI-Powered Supply Chain\nRisk Intelligence Assistant",
                 Inches(1.0), Inches(2.5), Inches(11.2), Inches(2.0),
                 size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Tagline
    add_text_box(sl,
        "Natural language queries  ·  RAG retrieval  ·  Multi-Agent AI  ·  Anomaly Detection",
        Inches(1.0), Inches(4.6), Inches(11), Inches(0.5),
        size=Pt(14), color=RGBColor(0xC8, 0xD3, 0xE8), align=PP_ALIGN.LEFT)

    # Bottom info row
    cols = [("Python 3.13", Inches(1.0)), ("FastAPI 0.136", Inches(3.2)),
            ("OpenAI GPT-4o-mini", Inches(5.4)), ("all-MiniLM + BM25 + NumPy", Inches(8.2))]
    for label, lx in cols:
        add_rect(sl, lx, Inches(5.5), Inches(2.0), Inches(0.42), fill=NAVY_LIGHT)
        add_text_box(sl, label, lx + Inches(0.1), Inches(5.55), Inches(1.8), Inches(0.35),
                     size=Pt(10), color=WHITE, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    sl = prs_slide(prs)
    page_header(sl, "Problem Statement",
                "The intelligence gap in global supply chain operations")

    pain_points = [
        ("Distributed Data Silos",
         "Operational data scattered across procurement, warehouse, logistics\nand vendor platforms — impossible to query holistically."),
        ("Reactive Monitoring Only",
         "Traditional systems trigger alerts only after thresholds breach.\nNo predictive capability for emerging disruption patterns."),
        ("Manual Analysis Bottleneck",
         "Operations teams spend hours consolidating reports to understand\nrisk — by the time insights arrive, disruptions have escalated."),
        ("Five Critical Scenarios",
         "Supplier delays  ·  Port congestion  ·  Stockout risk\nTransportation cost spikes  ·  Demand surge bottlenecks"),
    ]

    for i, (title, desc) in enumerate(pain_points):
        col = i % 2
        row = i // 2
        bx = Inches(0.45) + col * Inches(6.5)
        by = Inches(1.5) + row * Inches(2.6)
        bw = Inches(6.2)
        bh = Inches(2.3)
        card(sl, bx, by, bw, bh, bg=WHITE, border=ORANGE if i == 3 else NAVY)
        # Orange left bar
        add_rect(sl, bx, by, Inches(0.06), bh, fill=ORANGE if i < 3 else DANGER)
        add_text_box(sl, title, bx + Inches(0.18), by + Inches(0.15),
                     bw - Inches(0.3), Inches(0.38),
                     size=Pt(13), bold=True, color=NAVY)
        add_text_box(sl, desc, bx + Inches(0.18), by + Inches(0.55),
                     bw - Inches(0.3), Inches(1.5),
                     size=Pt(11), color=TEXT)


def slide_solution(prs):
    sl = prs_slide(prs)
    page_header(sl, "Our Solution",
                "AI-powered intelligence that transforms reactive monitoring into proactive resilience")

    # Central headline
    add_text_box(sl,
        "Operations teams can now query historical incidents in natural language,\n"
        "receive AI-generated risk assessments, and get explainable mitigation strategies\n"
        "— all within seconds.",
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.9),
        size=Pt(13.5), color=TEXT)

    caps = [
        ("🔍", "Natural Language\nIncident Retrieval", "Describe disruptions conversationally.\nSemantics + keywords."),
        ("⚖️", "4-Dimension\nRisk Scoring", "Supplier 35% · Inventory 30%\nShipment 25% · Demand 10%"),
        ("🤖", "Multi-Agent\nParallel Analysis", "4 specialist agents run concurrently\nwith A2A escalation."),
        ("🛡️", "Explainable\nRecommendations", "Every recommendation traces back\nto historical evidence."),
        ("📊", "Anomaly\nDetection", "IsolationForest flags statistical\nanomalies proactively."),
        ("🎯", "Evaluation\nFramework", "DeepEval RAG metrics +\nLLM-as-Judge quality scoring."),
    ]

    for i, (icon, title, desc) in enumerate(caps):
        col = i % 3
        row = i // 3
        bx = Inches(0.4) + col * Inches(4.3)
        by = Inches(2.6) + row * Inches(2.25)
        bw = Inches(4.0)
        bh = Inches(2.0)
        card(sl, bx, by, bw, bh, bg=CARD_BG, border=NAVY)
        add_text_box(sl, icon, bx + Inches(0.15), by + Inches(0.15),
                     Inches(0.55), Inches(0.55), size=Pt(22), align=PP_ALIGN.CENTER)
        add_text_box(sl, title, bx + Inches(0.75), by + Inches(0.15),
                     Inches(3.1), Inches(0.6), size=Pt(12), bold=True, color=NAVY)
        add_text_box(sl, desc, bx + Inches(0.15), by + Inches(0.8),
                     Inches(3.7), Inches(1.0), size=Pt(10.5), color=TEXT)


def slide_architecture(prs):
    """
    Matches the architecture diagram in TECHNICAL_DOCUMENT.md §2:
    FRONTEND → HTTP/REST → FASTAPI /api/v1 (4 routes)
      → Service Layer (RAG Search | RAG Generate | Orchestrator | Analytics)
      → CORE SERVICES (HybridSearch | RiskScoring | Guardrails | TokenOptimizer)
      → [NumPy VectorStore]   [OpenAI GPT-4o-mini via GW]
      → DATA LAYER (CSV → Preprocessing → Chunking → Embeddings)
    """
    sl = prs_slide(prs)
    page_header(sl, "System Architecture",
                "Layered microservice design (TECHNICAL_DOCUMENT.md §2) — each layer independently replaceable")

    COL_FG = RGBColor(0xC8, 0xD3, 0xE8)

    # ── Layer 1: FRONTEND ──────────────────────────────────────────────────
    add_rect(sl, Inches(0.4), Inches(1.52), Inches(12.5), Inches(0.58), fill=ORANGE)
    add_text_box(sl, "FRONTEND  (5-Tab SPA — Zero dependency HTML/CSS/JS)",
                 Inches(0.55), Inches(1.54), Inches(5.0), Inches(0.3),
                 size=Pt(10), bold=True, color=WHITE)
    tabs = ["Search", "Recommendations", "Multi-Agent", "Dashboard", "Anomalies"]
    for i, t in enumerate(tabs):
        tx = Inches(5.7) + i * Inches(1.42)
        add_rect(sl, tx, Inches(1.57), Inches(1.3), Inches(0.42),
                 fill=RGBColor(0xFF, 0xFF, 0xFF), line=None)
        add_text_box(sl, t, tx, Inches(1.57), Inches(1.3), Inches(0.42),
                     size=Pt(9), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Arrow
    add_text_box(sl, "▼  HTTP / REST", Inches(6.0), Inches(2.12),
                 Inches(2.5), Inches(0.28), size=Pt(9), bold=True,
                 color=ORANGE, align=PP_ALIGN.CENTER)

    # ── Layer 2: FASTAPI ────────────────────────────────────────────────────
    add_rect(sl, Inches(0.4), Inches(2.42), Inches(12.5), Inches(0.6), fill=NAVY)
    add_text_box(sl, "FASTAPI APPLICATION  /api/v1  (9 endpoints)",
                 Inches(0.55), Inches(2.44), Inches(4.5), Inches(0.3),
                 size=Pt(10), bold=True, color=ORANGE)
    routes = ["/search", "/recommendations", "/agents/analyze", "/analytics"]
    route_colors = [NAVY_LIGHT, NAVY_LIGHT, NAVY_LIGHT, NAVY_LIGHT]
    for i, r in enumerate(routes):
        rx = Inches(5.0) + i * Inches(1.97)
        add_rect(sl, rx, Inches(2.47), Inches(1.82), Inches(0.38),
                 fill=NAVY_LIGHT, line=ORANGE, line_w=Pt(0.8))
        add_text_box(sl, r, rx, Inches(2.47), Inches(1.82), Inches(0.38),
                     size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Layer 3: SERVICE LAYER ─────────────────────────────────────────────
    add_text_box(sl, "▼  Routes to Services", Inches(5.8), Inches(3.04),
                 Inches(3.0), Inches(0.24), size=Pt(8.5), bold=True,
                 color=COL_FG, align=PP_ALIGN.CENTER)

    svc_boxes = [
        ("RAG Service\n(search)", Inches(0.4), Inches(3.3), Inches(3.0), NAVY_LIGHT),
        ("RAG Service\n(generate+recommend)", Inches(3.55), Inches(3.3), Inches(3.0), NAVY_LIGHT),
        ("Orchestrator\n(multi-agent)", Inches(6.7), Inches(3.3), Inches(3.0), RGBColor(0x2E, 0x4A, 0x80)),
        ("Analytics Service\n(dashboard/anomalies)", Inches(9.85), Inches(3.3), Inches(3.05), RGBColor(0x2E, 0x4A, 0x80)),
    ]
    for label, sx, sy, sw, bg in svc_boxes:
        add_rect(sl, sx, sy, sw, Inches(0.65), fill=bg)
        add_text_box(sl, label, sx, sy, sw, Inches(0.65),
                     size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Layer 4: CORE SERVICES ─────────────────────────────────────────────
    add_text_box(sl, "▼  Shared Core Services", Inches(5.5), Inches(3.97),
                 Inches(3.0), Inches(0.24), size=Pt(8.5), bold=True,
                 color=COL_FG, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.4), Inches(4.23), Inches(12.5), Inches(0.62),
             fill=RGBColor(0x18, 0x2E, 0x60))
    core_items = [
        ("HybridSearch\n(BM25 + RRF)", Inches(0.5)),
        ("RiskScoring\n(4-dimension)", Inches(3.7)),
        ("Guardrails\n(validation)", Inches(6.9)),
        ("TokenOptimizer\n(tiktoken)", Inches(10.1)),
    ]
    add_text_box(sl, "CORE SERVICES LAYER",
                 Inches(0.55), Inches(4.25), Inches(2.8), Inches(0.28),
                 size=Pt(8.5), bold=True, color=ORANGE)
    for label, lx in core_items:
        add_rect(sl, lx, Inches(4.27), Inches(3.0), Inches(0.52),
                 fill=RGBColor(0x24, 0x3D, 0x72), line=ORANGE, line_w=Pt(0.7))
        add_text_box(sl, label, lx, Inches(4.27), Inches(3.0), Inches(0.52),
                     size=Pt(9), color=WHITE, align=PP_ALIGN.CENTER)

    # ── Layer 5: VectorStore (left) + LLM (right) ──────────────────────────
    # Left: NumPy VectorStore
    add_rect(sl, Inches(0.4), Inches(4.95), Inches(5.8), Inches(1.08),
             fill=RGBColor(0x12, 0x22, 0x48), line=NAVY_LIGHT, line_w=Pt(1))
    add_text_box(sl, "NUMPY VECTOR STORE",
                 Inches(0.55), Inches(4.98), Inches(3.5), Inches(0.3),
                 size=Pt(9.5), bold=True, color=ORANGE)
    add_text_box(sl, "embeddings.npy  ·  documents.json  ·  BM25 Index\n"
                     "all-MiniLM-L6-v2 (384-dim) · store_meta.json (dimension guard)",
                 Inches(0.55), Inches(5.3), Inches(5.5), Inches(0.65),
                 size=Pt(9), color=COL_FG)

    # Right: OpenAI GPT-4o-mini
    add_rect(sl, Inches(7.1), Inches(4.95), Inches(5.8), Inches(1.08),
             fill=ORANGE, line=WHITE, line_w=Pt(1.2))
    add_text_box(sl, "OPENAI GPT-4o-mini  (via keygateway · SSL-off)",
                 Inches(7.25), Inches(4.98), Inches(5.5), Inches(0.3),
                 size=Pt(9.5), bold=True, color=WHITE)
    add_text_box(sl, "• Recommendation generation (chat.completions, max_tokens=500)\n"
                     "• Multi-agent tool use (function calling)  ·  LLM-as-judge evaluation",
                 Inches(7.25), Inches(5.3), Inches(5.5), Inches(0.65),
                 size=Pt(9), color=RGBColor(0xFF, 0xE8, 0xD8))

    # ── Layer 6: DATA LAYER ────────────────────────────────────────────────
    add_rect(sl, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.65),
             fill=RGBColor(0x0A, 0x16, 0x30))
    add_text_box(sl, "DATA LAYER",
                 Inches(0.55), Inches(6.18), Inches(1.8), Inches(0.28),
                 size=Pt(9.5), bold=True, color=ORANGE)
    add_text_box(sl,
                 "supply_chain_data.csv (600 incidents)  →  Preprocessing (coerce/enrich)  →  "
                 "Chunking (tiktoken)  →  Embeddings  →  VectorStore upsert",
                 Inches(2.5), Inches(6.18), Inches(10.3), Inches(0.6),
                 size=Pt(9.5), color=COL_FG)


def slide_data_flow(prs):
    sl = prs_slide(prs)
    page_header(sl, "Data Flow Diagram",
                "From raw CSV incidents to explainable AI recommendations")

    # Ingestion path (top row)
    add_text_box(sl, "DATA INGESTION FLOW", Inches(0.5), Inches(1.55),
                 Inches(12), Inches(0.3), size=Pt(10), bold=True,
                 color=ORANGE, align=PP_ALIGN.LEFT)

    ing_steps = [
        ("CSV\n600 Incidents", NAVY),
        ("Preprocess\n& Enrich", NAVY_LIGHT),
        ("Chunk\n+ Tokenize", NAVY_LIGHT),
        ("Encode\nEmbeddings\n384-dim", ORANGE),
        ("NumPy\nVectorStore", NAVY),
        ("BM25\nIndex", NAVY),
    ]
    step_w = Inches(1.9)
    step_h = Inches(0.9)
    sy = Inches(1.9)
    for i, (label, bg) in enumerate(ing_steps):
        sx = Inches(0.4) + i * Inches(2.15)
        add_rect(sl, sx, sy, step_w, step_h, fill=bg)
        add_text_box(sl, label, sx, sy, step_w, step_h,
                     size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(ing_steps) - 1:
            add_text_box(sl, "→", sx + step_w, sy + Inches(0.25),
                         Inches(0.25), Inches(0.4),
                         size=Pt(16), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(sl, Inches(0.4), Inches(2.95), Inches(12.5), Inches(0.03), fill=NAVY_LIGHT)

    # Query path (bottom rows)
    add_text_box(sl, "QUERY & RESPONSE FLOW", Inches(0.5), Inches(3.05),
                 Inches(12), Inches(0.3), size=Pt(10), bold=True,
                 color=ORANGE, align=PP_ALIGN.LEFT)

    query_steps = [
        ("User\nQuery", NAVY),
        ("Guardrails\nValidate", RGBColor(0xC4, 0x7A, 0x00)),
        ("Hybrid\nSearch\nBM25+Sem", NAVY_LIGHT),
        ("RRF\nFusion\n+ Rerank", NAVY_LIGHT),
        ("Token\nBudget\nOptimize", NAVY_LIGHT),
        ("Risk\nScore\n4-Dim", ORANGE),
        ("GPT-4o-mini\nGenerate", ORANGE),
        ("Parse &\nReturn\nResponse", NAVY),
    ]
    step_w2 = Inches(1.5)
    sy2 = Inches(3.38)
    for i, (label, bg) in enumerate(query_steps):
        sx = Inches(0.35) + i * Inches(1.63)
        add_rect(sl, sx, sy2, step_w2, Inches(0.9), fill=bg)
        add_text_box(sl, label, sx, sy2, step_w2, Inches(0.9),
                     size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(query_steps) - 1:
            add_text_box(sl, "→", sx + step_w2, sy2 + Inches(0.25),
                         Inches(0.13), Inches(0.4),
                         size=Pt(13), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Multi-agent parallel
    add_text_box(sl, "MULTI-AGENT PARALLEL FLOW  (A2A Escalation)", Inches(0.5), Inches(4.5),
                 Inches(12), Inches(0.3), size=Pt(10), bold=True, color=ORANGE)

    agents = [
        ("Supplier\nRisk Agent", NAVY),
        ("Shipment\nAnalysis Agent", NAVY_LIGHT),
        ("Inventory\nIntelligence Agent", RGBColor(0x2E, 0x4A, 0x80)),
    ]
    for i, (label, bg) in enumerate(agents):
        ax = Inches(0.4) + i * Inches(3.5)
        add_rect(sl, ax, Inches(4.85), Inches(3.2), Inches(0.8), fill=bg)
        add_text_box(sl, label, ax, Inches(4.85), Inches(3.2), Inches(0.8),
                     size=Pt(10.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(sl, "↓  ↓  ↓  A2A Escalation  ↓  ↓  ↓",
                 Inches(0.4), Inches(5.72), Inches(10.2), Inches(0.3),
                 size=Pt(11), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(3.0), Inches(6.08), Inches(4.5), Inches(0.75), fill=ORANGE)
    add_text_box(sl, "Recommendation Agent  →  Consolidated Response",
                 Inches(3.0), Inches(6.08), Inches(4.5), Inches(0.75),
                 size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_rag_pipeline(prs):
    sl = prs_slide(prs)
    page_header(sl, "RAG Pipeline",
                "8-stage Retrieve-Augment-Generate pipeline with hybrid search")

    steps = [
        ("1", "User Query\nReceived",          "Natural language input\nfrom operations team",           NAVY),
        ("2", "Guardrails\nValidation",         "Injection detection\nDomain relevance check",            RGBColor(0xC8, 0x20, 0x1A)),
        ("3", "Hybrid Search\nBM25 + Semantic", "Parallel BM25 + cosine ANN\nMerged via RRF (α=0.5)",    NAVY_LIGHT),
        ("4", "Logistics\nReranking",           "+0.05 boost per matching\ndomain field",                 RGBColor(0x2E, 0x4A, 0x80)),
        ("5", "Token Budget\nManagement",       "tiktoken cl100k_base\nmax 8000 token context",           RGBColor(0x36, 0x56, 0x98)),
        ("6", "4-Dim Risk\nScoring",            "Supplier · Inventory\nShipment · Demand",                ORANGE),
        ("7", "GPT-4o-mini\nGeneration",        "Structured prompt injection\nPRIORITY 1-5 output format",ORANGE),
        ("8", "Parse &\nReturn",                "Regex extraction\nFallback to plain text",               NAVY),
    ]

    for i, (num, title, detail, bg) in enumerate(steps):
        col = i % 4
        row = i // 4
        bx = Inches(0.35) + col * Inches(3.25)
        by = Inches(1.58) + row * Inches(2.6)
        bw = Inches(3.0)
        bh = Inches(2.3)
        # Number badge
        add_rect(sl, bx, by, Inches(0.45), Inches(0.45), fill=ORANGE)
        add_text_box(sl, num, bx, by, Inches(0.45), Inches(0.45),
                     size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Card
        add_rect(sl, bx, by + Inches(0.45), bw, bh - Inches(0.45), fill=bg)
        add_text_box(sl, title, bx + Inches(0.1), by + Inches(0.55),
                     bw - Inches(0.2), Inches(0.7),
                     size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, detail, bx + Inches(0.1), by + Inches(1.3),
                     bw - Inches(0.2), Inches(1.0),
                     size=Pt(10), color=RGBColor(0xC8, 0xD3, 0xE8), align=PP_ALIGN.CENTER)
        # Arrow
        if col < 3:
            add_text_box(sl, "→", bx + bw + Inches(0.05), by + Inches(1.1),
                         Inches(0.2), Inches(0.4),
                         size=Pt(16), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_multiagent(prs):
    sl = prs_slide(prs)
    page_header(sl, "Multi-Agent Architecture",
                "Parallel specialist agents with Agent-to-Agent (A2A) escalation workflow")

    # Orchestrator box
    add_rect(sl, Inches(4.9), Inches(1.55), Inches(3.5), Inches(0.72), fill=NAVY)
    add_text_box(sl, "🎯  Multi-Agent Orchestrator",
                 Inches(4.9), Inches(1.55), Inches(3.5), Inches(0.72),
                 size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Context retrieval
    add_rect(sl, Inches(0.4), Inches(1.55), Inches(4.1), Inches(0.72), fill=NAVY_LIGHT)
    add_text_box(sl, "Hybrid Search  →  Context Docs (top 10)",
                 Inches(0.4), Inches(1.55), Inches(4.1), Inches(0.72),
                 size=Pt(10.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow down from orchestrator
    add_text_box(sl, "↓  ThreadPoolExecutor (max_workers=3)  ↓",
                 Inches(3.0), Inches(2.4), Inches(7.0), Inches(0.38),
                 size=Pt(11), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # 3 specialist agents
    agent_data = [
        ("Supplier\nRisk Agent",
         "get_supplier_history()\nDelay · SLA · Dependency\nEscalate if risk > 0.7",
         NAVY),
        ("Shipment\nAnalysis Agent",
         "get_route_status()\nPort · Carrier · Customs\nRoot cause analysis",
         NAVY_LIGHT),
        ("Inventory\nIntelligence Agent",
         "check_inventory()\nget_demand_forecast()\nStockout prediction",
         RGBColor(0x2E, 0x4A, 0x80)),
    ]
    for i, (title, detail, bg) in enumerate(agent_data):
        ax = Inches(0.4) + i * Inches(4.3)
        add_rect(sl, ax, Inches(2.85), Inches(4.0), Inches(1.85), fill=bg)
        add_text_box(sl, title, ax + Inches(0.1), Inches(2.9),
                     Inches(3.8), Inches(0.6),
                     size=Pt(12.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, detail, ax + Inches(0.1), Inches(3.5),
                     Inches(3.8), Inches(1.1),
                     size=Pt(9.5), color=RGBColor(0xC8, 0xD3, 0xE8), align=PP_ALIGN.CENTER)

    # A2A Escalation
    add_text_box(sl, "⚡  A2A Escalation Chain  (risk > 0.7 → RecommendationAgent → Director if > 0.85)",
                 Inches(0.4), Inches(4.82), Inches(12.5), Inches(0.38),
                 size=Pt(11), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Recommendation agent
    add_rect(sl, Inches(3.5), Inches(5.3), Inches(6.3), Inches(0.85), fill=ORANGE)
    add_text_box(sl, "🎯  Recommendation Agent  —  Synthesise · Consolidate · Proactive Alerts",
                 Inches(3.5), Inches(5.3), Inches(6.3), Inches(0.85),
                 size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Output
    add_rect(sl, Inches(2.8), Inches(6.25), Inches(7.7), Inches(0.72), fill=NAVY)
    add_text_box(sl, "OrchestratorResponse  ·  consolidated_risk  ·  escalation_chain  ·  proactive_alerts",
                 Inches(2.8), Inches(6.25), Inches(7.7), Inches(0.72),
                 size=Pt(10.5), color=RGBColor(0xC8, 0xD3, 0xE8), align=PP_ALIGN.CENTER)


def slide_tech_stack(prs):
    sl = prs_slide(prs)
    page_header(sl, "Technology Stack",
                "Production-grade open-source stack running on Python 3.13")

    categories = [
        ("Backend Framework",
         [("FastAPI 0.136", "REST API + OpenAPI docs"), ("Uvicorn", "ASGI server"),
          ("Pydantic v2", "Request/response validation")],
         NAVY),
        ("AI / LLM",
         [("OpenAI GPT-4o-mini", "Chat completions via gateway"), ("OpenAI API SDK 1.55", "Function calling"),
          ("DeepEval (optional)", "RAG quality metrics")],
         ORANGE),
        ("Retrieval & Search",
         [("all-MiniLM-L6-v2", "384-dim sentence embeddings"), ("BM25Okapi", "Sparse keyword retrieval"),
          ("RRF Fusion", "Reciprocal Rank Fusion")],
         NAVY_LIGHT),
        ("Vector Store & Data",
         [("NumPy + JSON", "Pure-Python vector store (no MSVC)"), ("tiktoken cl100k_base", "Token budget mgmt"),
          ("Pandas + scikit-learn", "Data processing + IsolationForest")],
         RGBColor(0x2E, 0x4A, 0x80)),
    ]

    for i, (cat, items, bg) in enumerate(categories):
        col = i % 2
        row = i // 2
        bx = Inches(0.4) + col * Inches(6.5)
        by = Inches(1.55) + row * Inches(2.55)
        bw = Inches(6.1)
        bh = Inches(2.3)
        add_rect(sl, bx, by, bw, Inches(0.45), fill=bg)
        add_text_box(sl, cat, bx + Inches(0.1), by + Inches(0.06),
                     bw - Inches(0.2), Inches(0.38),
                     size=Pt(12), bold=True, color=WHITE)
        add_rect(sl, bx, by + Inches(0.45), bw, bh - Inches(0.45),
                 fill=CARD_BG, line=bg, line_w=Pt(1))
        for j, (tech, desc) in enumerate(items):
            ty = by + Inches(0.55) + j * Inches(0.55)
            add_rect(sl, bx + Inches(0.1), ty + Inches(0.04),
                     Inches(0.06), Inches(0.3), fill=ORANGE)
            add_text_box(sl, f"{tech}  —  {desc}",
                         bx + Inches(0.24), ty, bw - Inches(0.35), Inches(0.45),
                         size=Pt(11), color=TEXT)


def slide_evaluation(prs):
    sl = prs_slide(prs)
    page_header(sl, "Evaluation Framework",
                "Dual validation: DeepEval RAG metrics + custom LLM-as-Judge")

    # DeepEval column
    add_rect(sl, Inches(0.4), Inches(1.55), Inches(5.9), Inches(0.5), fill=NAVY)
    add_text_box(sl, "DeepEval — RAG Quality Metrics",
                 Inches(0.4), Inches(1.55), Inches(5.9), Inches(0.5),
                 size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    metrics = [
        ("Answer Relevancy", "Does the response answer the query?", "≥ 0.5"),
        ("Faithfulness", "Are claims grounded in retrieved context? (anti-hallucination)", "≥ 0.5"),
        ("Contextual Precision", "Are retrieved documents relevant to the question?", "≥ 0.5"),
        ("Contextual Recall", "Does context cover all necessary information?", "≥ 0.5"),
    ]
    for i, (m, desc, threshold) in enumerate(metrics):
        by = Inches(2.15) + i * Inches(0.72)
        add_rect(sl, Inches(0.4), by, Inches(5.9), Inches(0.65), fill=CARD_BG, line=NAVY, line_w=Pt(0.8))
        add_text_box(sl, m, Inches(0.55), by + Inches(0.08), Inches(2.4), Inches(0.38),
                     size=Pt(11), bold=True, color=NAVY)
        add_text_box(sl, desc, Inches(2.9), by + Inches(0.08), Inches(2.5), Inches(0.38),
                     size=Pt(9.5), color=TEXT)
        add_rect(sl, Inches(5.6), by + Inches(0.1), Inches(0.65), Inches(0.42), fill=ORANGE)
        add_text_box(sl, threshold, Inches(5.6), by + Inches(0.1), Inches(0.65), Inches(0.42),
                     size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # LLM Judge column
    add_rect(sl, Inches(6.8), Inches(1.55), Inches(6.1), Inches(0.5), fill=ORANGE)
    add_text_box(sl, "LLM-as-Judge — Business Quality Scoring",
                 Inches(6.8), Inches(1.55), Inches(6.1), Inches(0.5),
                 size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    judge_dims = [
        ("Relevance", "Are recs tied to the identified risk?"),
        ("Actionability", "Can teams execute these immediately?"),
        ("Completeness", "All risk dimensions addressed?"),
        ("Feasibility", "Realistic timelines and owners?"),
        ("Evidence-Based", "Grounded in historical incidents?"),
    ]
    for i, (dim, desc) in enumerate(judge_dims):
        by = Inches(2.15) + i * Inches(0.72)
        add_rect(sl, Inches(6.8), by, Inches(6.1), Inches(0.65),
                 fill=CARD_BG, line=ORANGE, line_w=Pt(0.8))
        add_text_box(sl, dim, Inches(6.95), by + Inches(0.08), Inches(1.9), Inches(0.38),
                     size=Pt(11), bold=True, color=ORANGE)
        add_text_box(sl, desc, Inches(8.9), by + Inches(0.08), Inches(3.8), Inches(0.38),
                     size=Pt(9.5), color=TEXT)
        add_text_box(sl, "0–10", Inches(12.3), by + Inches(0.15), Inches(0.5), Inches(0.35),
                     size=Pt(9), bold=True, color=NAVY)

    # Verdict chips
    add_text_box(sl, "Verdict:", Inches(6.8), Inches(5.8), Inches(1.5), Inches(0.4),
                 size=Pt(11), bold=True, color=NAVY)
    for label, bg, tx in [("APPROVED", SUCCESS, WHITE),
                           ("NEEDS IMPROVEMENT", WARNING, TEXT),
                           ("REJECTED", DANGER, WHITE)]:
        lx = Inches(6.8) + [0, 2.1, 4.7][["APPROVED", "NEEDS IMPROVEMENT", "REJECTED"].index(label)]
        add_rect(sl, lx, Inches(5.78), Inches(1.95 if label == "NEEDS IMPROVEMENT" else 1.8), Inches(0.42), fill=bg)
        add_text_box(sl, label, lx, Inches(5.78), Inches(2.0), Inches(0.42),
                     size=Pt(9.5), bold=True, color=tx, align=PP_ALIGN.CENTER)


def slide_fallbacks(prs):
    sl = prs_slide(prs)
    page_header(sl, "Fallback & Resilience Mechanisms",
                "9 graceful degradation patterns — system never fails silently")

    fallbacks = [
        ("Embedding Model",
         "OpenAI /embeddings unreachable",
         "Local all-MiniLM-L6-v2 (384-dim, no API)", ORANGE),
        ("Vector Store",
         "Model or dimension changed",
         "Auto-clear stale index + re-ingest on startup", NAVY),
        ("LLM Parsing",
         "No PRIORITY blocks in GPT response",
         "First 5 lines wrapped as generic steps", NAVY_LIGHT),
        ("DeepEval",
         "Library not installed / API error",
         "Placeholder scores (0.70–0.80) with note", RGBColor(0x2E, 0x4A, 0x80)),
        ("Agent Crash",
         "Agent timeout (60s) or exception",
         "Stub AgentResult (risk=0.0, manual review)", RGBColor(0xC8, 0x20, 0x1A)),
        ("Agent Routing",
         "No keywords match query",
         "Run all 3 specialist agents (safe default)", NAVY),
        ("BM25 Index",
         "Index not yet built",
         "Auto-rebuild from vector store corpus", NAVY_LIGHT),
        ("Server Startup",
         "Dataset CSV not found",
         "Server starts with empty store (search → [])", RGBColor(0xC4, 0x7A, 0x00)),
        ("Gateway Down",
         "LLM chat completions timeout",
         "Returns 400 with clear actionable message", RGBColor(0xC8, 0x20, 0x1A)),
    ]

    for i, (layer, trigger, fallback, bg) in enumerate(fallbacks):
        col = i % 3
        row = i // 3
        bx = Inches(0.4) + col * Inches(4.3)
        by = Inches(1.58) + row * Inches(1.8)
        bw = Inches(4.0)
        add_rect(sl, bx, by, bw, Inches(0.38), fill=bg)
        add_text_box(sl, layer, bx + Inches(0.1), by + Inches(0.04), bw - Inches(0.2), Inches(0.3),
                     size=Pt(10.5), bold=True, color=WHITE)
        add_rect(sl, bx, by + Inches(0.38), bw, Inches(1.3),
                 fill=CARD_BG, line=bg, line_w=Pt(0.8))
        add_text_box(sl, f"Trigger: {trigger}", bx + Inches(0.1),
                     by + Inches(0.45), bw - Inches(0.2), Inches(0.45),
                     size=Pt(9.5), color=MUTED, italic=True)
        add_text_box(sl, f"→  {fallback}", bx + Inches(0.1),
                     by + Inches(0.9), bw - Inches(0.2), Inches(0.7),
                     size=Pt(9.5), bold=True, color=TEXT)


def slide_dashboard(prs):
    sl = prs_slide(prs)
    page_header(sl, "5-Tab Frontend Dashboard",
                "Zero-dependency SPA — no build tools, no npm, opens directly from filesystem")

    tabs = [
        ("🔍  Search",
         "Hybrid BM25 + Semantic\nnatural language incident search\n\nFilters: Severity · Status · Supplier\nLatency: 10–40 ms (local embeddings)", NAVY),
        ("💡  Recommendations",
         "AI-generated mitigation strategies\nRisk bars · PRIORITY 1-5 steps\n\nTimeline · Owner · Expected Impact\nOptional LLM Judge evaluation", NAVY_LIGHT),
        ("🤖  Multi-Agent",
         "Parallel agent analysis panel\nA2A escalation chain display\n\nProactive alerts · Consolidated risk\nAll 4 agents visible with scores", RGBColor(0x2E, 0x4A, 0x80)),
        ("📊  Dashboard",
         "Live KPI analytics grid\n600 incidents · 302 critical\n\nHigh-risk suppliers · Stockout locs\nDisruption trend (improving/stable/bad)", ORANGE),
        ("⚠️  Anomalies",
         "IsolationForest anomaly detection\nZ-score feature attribution\n\nPearson correlation insights\nContamination rate configurable", RGBColor(0xC8, 0x20, 0x1A)),
    ]

    for i, (title, desc, bg) in enumerate(tabs):
        tx = Inches(0.4) + i * Inches(2.56)
        add_rect(sl, tx, Inches(1.55), Inches(2.4), Inches(0.52), fill=bg)
        add_text_box(sl, title, tx, Inches(1.55), Inches(2.4), Inches(0.52),
                     size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, tx, Inches(2.07), Inches(2.4), Inches(4.85),
                 fill=CARD_BG, line=bg, line_w=Pt(1.2))
        add_text_box(sl, desc, tx + Inches(0.12), Inches(2.18),
                     Inches(2.18), Inches(4.6), size=Pt(10.5), color=TEXT)

    # Sample queries note
    add_text_box(sl, "✨  6 clickable sample query chips per tab for easy exploration",
                 Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
                 size=Pt(11), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_demo_flow(prs):
    sl = prs_slide(prs)
    page_header(sl, "Demo Flow  —  10 Minutes",
                "Operations manager investigates a supply chain disruption end-to-end")

    steps = [
        ("00:00\n–01:00", "Problem\nStatement",
         "Explain the supply chain intelligence gap.\nDistributed data · Reactive monitoring · Manual analysis.", NAVY),
        ("01:00\n–02:00", "System\nOverview",
         "Walk through architecture slide.\nFrontend → API → RAG → Agents → LLM.", NAVY),
        ("02:00\n–03:30", "Search\nTab",
         "Query: 'supplier delivery delays for critical components'\nShow hybrid results, severity chips, latency 40ms.", NAVY_LIGHT),
        ("03:30\n–05:30", "Recommendations\nTab",
         "Query: 'port congestion impacting shipment schedules'\nHighlight animated risk bars + 5 priority mitigations.", ORANGE),
        ("05:30\n–07:30", "Multi-Agent\nTab",
         "Query: 'warehouse inventory approaching stockout'\nShow 4 agent cards · A2A escalation chain · Alerts.", RGBColor(0x2E, 0x4A, 0x80)),
        ("07:30\n–08:30", "Dashboard &\nAnomalies",
         "KPI grid: 600 incidents · 302 critical · Deteriorating.\nAnomaly cards with Z-score attribution.", RGBColor(0xC8, 0x20, 0x1A)),
        ("08:30\n–10:00", "Q&A",
         "Panel questions.\nKey differentiators:\nRRF · 4-dim risk · A2A · Graceful fallbacks.", SUCCESS),
    ]

    step_w = Inches(1.68)
    for i, (time, title, detail, bg) in enumerate(steps):
        bx = Inches(0.35) + i * Inches(1.84)
        add_rect(sl, bx, Inches(1.58), step_w, Inches(0.58), fill=bg)
        add_text_box(sl, time, bx, Inches(1.58), step_w, Inches(0.58),
                     size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, bx, Inches(2.16), step_w, Inches(0.55), fill=ORANGE)
        add_text_box(sl, title, bx, Inches(2.16), step_w, Inches(0.55),
                     size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, bx, Inches(2.71), step_w, Inches(4.45),
                 fill=CARD_BG, line=bg, line_w=Pt(1))
        add_text_box(sl, detail, bx + Inches(0.08), Inches(2.82),
                     step_w - Inches(0.16), Inches(4.2), size=Pt(9.5), color=TEXT)
        if i < len(steps) - 1:
            add_text_box(sl, "→", bx + step_w + Inches(0.05), Inches(2.85),
                         Inches(0.16), Inches(0.4),
                         size=Pt(13), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_results(prs):
    sl = prs_slide(prs)
    page_header(sl, "Performance & Results",
                "Live test results on 600-incident corpus with local sentence-transformer embeddings")

    kpis = [
        ("40 ms",   "Search Latency",    "Hybrid BM25+Semantic (local embeddings)", NAVY),
        ("600",     "Indexed Incidents", "2024–2026 synthetic supply chain records", NAVY_LIGHT),
        ("30",      "Anomalies Found",   "IsolationForest contamination=0.05",       ORANGE),
        ("9.19d",   "Avg Delivery Delay","Across all incidents (deteriorating trend)",RGBColor(0xC4, 0x7A, 0x00)),
        ("302",     "Critical Incidents","50.3% of corpus flagged as critical severity",RGBColor(0xC8, 0x20, 0x1A)),
        ("4 Agents","Multi-Agent",       "Parallel + A2A escalation in ~22s",        SUCCESS),
    ]

    for i, (val, label, note, bg) in enumerate(kpis):
        col = i % 3
        row = i // 3
        bx = Inches(0.4) + col * Inches(4.3)
        by = Inches(1.58) + row * Inches(2.55)
        bw = Inches(4.0)
        add_rect(sl, bx, by, bw, Inches(0.06), fill=bg)
        add_rect(sl, bx, by + Inches(0.06), bw, Inches(2.25),
                 fill=CARD_BG, line=NAVY, line_w=Pt(0.8))
        add_text_box(sl, val, bx, by + Inches(0.18),
                     bw, Inches(0.9), size=Pt(32), bold=True, color=bg, align=PP_ALIGN.CENTER)
        add_text_box(sl, label, bx, by + Inches(1.1),
                     bw, Inches(0.4), size=Pt(13), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(sl, note, bx + Inches(0.1), by + Inches(1.55),
                     bw - Inches(0.2), Inches(0.7), size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)


def slide_thankyou(prs):
    sl = prs_slide(prs)
    set_bg(sl, NAVY_BG)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=ORANGE)
    add_rect(sl, 0, Inches(7.38), W, Inches(0.12), fill=ORANGE)
    add_rect(sl, 0, Inches(0.12), Inches(0.1), Inches(7.26), fill=ORANGE)

    add_text_box(sl, "⬡",
                 Inches(5.9), Inches(1.3), Inches(1.5), Inches(1.3),
                 size=Pt(60), color=ORANGE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(3.0), Inches(2.7), Inches(7.3), Inches(0.04), fill=ORANGE)
    add_text_box(sl, "Thank You",
                 Inches(2.0), Inches(2.85), Inches(9.3), Inches(1.1),
                 size=Pt(44), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, "AI-Powered Supply Chain Risk Intelligence Assistant",
                 Inches(2.5), Inches(4.0), Inches(8.3), Inches(0.55),
                 size=Pt(16), color=RGBColor(0xC8, 0xD3, 0xE8), align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(3.0), Inches(4.65), Inches(7.3), Inches(0.04), fill=ORANGE)

    add_text_box(sl, "Siddharth SR  |  siddharth.sr@prodapt.com  |  2026",
                 Inches(2.5), Inches(4.85), Inches(8.3), Inches(0.45),
                 size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

    qs = [
        ("🔍 Search", "Hybrid BM25+Semantic"),
        ("🤖 Agents", "A2A Escalation"),
        ("📊 Analytics", "IsolationForest"),
        ("💡 RAG", "GPT-4o-mini"),
        ("🛡️ Eval", "LLM-as-Judge"),
    ]
    for i, (title, sub) in enumerate(qs):
        qx = Inches(1.3) + i * Inches(2.2)
        add_rect(sl, qx, Inches(5.7), Inches(1.9), Inches(1.35), fill=NAVY_LIGHT)
        add_text_box(sl, title, qx, Inches(5.78), Inches(1.9), Inches(0.48),
                     size=Pt(12), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_text_box(sl, sub, qx, Inches(6.26), Inches(1.9), Inches(0.38),
                     size=Pt(9.5), color=RGBColor(0xC8, 0xD3, 0xE8), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════
def build_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    slide_title(prs);        print("  1. Title")
    slide_problem(prs);      print("  2. Problem Statement")
    slide_solution(prs);     print("  3. Solution Overview")
    slide_architecture(prs); print("  4. System Architecture")
    slide_data_flow(prs);    print("  5. Data Flow Diagram")
    slide_rag_pipeline(prs); print("  6. RAG Pipeline")
    slide_multiagent(prs);   print("  7. Multi-Agent Architecture")
    slide_tech_stack(prs);   print("  8. Technology Stack")
    slide_evaluation(prs);   print("  9. Evaluation Framework")
    slide_fallbacks(prs);    print(" 10. Fallback Mechanisms")
    slide_dashboard(prs);    print(" 11. Frontend Dashboard")
    slide_results(prs);      print(" 12. Results & Performance")
    slide_thankyou(prs);     print(" 13. Thank You / Q&A")

    out = Path(__file__).parent / "SC_Risk_Intelligence_PPT.pptx"
    prs.save(str(out))
    print(f"\nSaved -> {out}")
    return str(out)


if __name__ == "__main__":
    build_presentation()
